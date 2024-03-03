import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from langchain.text_splitter import CharacterTextSplitter
from langchain_community.embeddings import HuggingFaceInstructEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_community.embeddings import SentenceTransformerEmbeddings
from langchain.memory import ConversationBufferMemory
from langchain.chains import ConversationalRetrievalChain
from langchain.llms import HuggingFaceHub
from htmlTemplates import css, bot_template, user_template

# def get_pdf_text(pdf_docs):
#     text = ""
#     for pdf in pdf_docs:
#         if pdf.name.endswith('.pdf'):
#             pdf_reader = PdfReader(pdf)
#             for page in pdf_reader.pages:
#                 text += page.extract_text()
#         elif pdf.name.endswith('.pptx'):
#             prs = Presentation(pdf)
#             for slide in prs.slides:
#                 for shape in slide.shapes:
#                     if hasattr(shape, "text"):
#                         text += shape.text
#         elif pdf.name.endswith('.docx'):
#             doc = Document(pdf)
#             for paragraph in doc.paragraphs:
#                 text += paragraph.text + "\n"
#     return text


def get_text_from_pdf(pdf_file):
    pdf_reader = PdfReader(pdf_file)
    text = ""
    for page in pdf_reader.pages:
        text += page.extract_text()
    return text

def get_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
    return text

def get_text_from_docx(docx_file):
    doc = Document(docx_file)
    text = ""
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

# def get_text_from_txt(txt_file):
#     file_contents = txt_file
#     print(file_contents.decode('utf-8'))
        
#     return file_contents

def get_text_chunks(text):
    text_splitter = CharacterTextSplitter(
        separator="\n",
        chunk_size = 1000, 
        chunk_overlap = 200,
        length_function = len
    )
    
    chunks = text_splitter.split_text(text)
    return chunks

def get_vectorstore(text_chunks):
    embeddings = SentenceTransformerEmbeddings(model_name="hkunlp/instructor-xl")
    # embeddings = HuggingFaceInstructEmbeddings(model_name="hkunlp/instructor-xl")
    vectorstore = FAISS.from_texts(texts = text_chunks, embedding = embeddings)
    return vectorstore

def get_conversation_chain(vectorstore):
    # llm = ChatOpenAI()
    
    llm = HuggingFaceHub(repo_id="google/flan-t5-xxl", model_kwargs={"temperature":0.5, "max_length":512})
    memory = ConversationBufferMemory(memory_key='chat_history',return_messages=True)
    conversation_chain = ConversationalRetrievalChain.from_llm(
        llm= llm,
        retriever=vectorstore.as_retriever(),
        memory = memory
    )
    return conversation_chain

def handle_userinput(user_question):
    response = st.session_state.conversation({'question': user_question})
    st.session_state.chat_history = response['chat_history']
    for i, message in enumerate(st.session_state.chat_history):
        if i & 2 ==0:
            st.write(user_template.replace("{{MSG}}",message.content),unsafe_allow_html=True)
        else:
            st.write(bot_template.replace("{{MSG}}",message.content),unsafe_allow_html=True)
            

def main():
    load_dotenv()
    st.set_page_config(page_title="chat with multiple pdf", page_icon=":books:")
    
    st.write(css, unsafe_allow_html=True)
    if "conversation" not in st.session_state:
        st.session_state.conversation = None

    if "chat_history" not in st.session_state:
        st.session_state.chat_history = None

    st.header("chat with multiple PDFs :books:")
    user_question = st.text_input("Ask a question about your document:")
    if user_question:
        handle_userinput(user_question)
    # st.write(user_template.replace("{{MSG}}","Hello Robot"),unsafe_allow_html=True) 
    # st.write(bot_template.replace("{{MSG}}","Hello Human"),unsafe_allow_html=True)
    


    with st.sidebar:
        st.subheader("your documents")
        pdf_docs = st.file_uploader("Upload your pdfs here and click on process", accept_multiple_files=True)
        # print(pdf_docs[1])
        if st.button("process"):
            with st.spinner("processing"):
                # get pdf text
                raw_text = ""
                for file in pdf_docs:
                    
                    file_type = file.name.split('.')[-1]
                    if file_type == 'pdf':
                        raw_text += get_text_from_pdf(file)
                    elif file_type == 'pptx':
                        raw_text += get_text_from_pptx(file)
                    elif file_type == 'docx':
                        raw_text += get_text_from_docx(file)
                    # elif file_type == 'txt':
                    #     raw_text += get_text_from_txt(file)
                # raw_text = get_pdf_text(pdf_docs)   
         
                # get the text chunks 
                text_chunks = get_text_chunks(raw_text)

                # create vector store
                vectorstore = get_vectorstore(text_chunks)

                st.session_state.conversation = get_conversation_chain(vectorstore)



if __name__ == '__main__':
    main()
